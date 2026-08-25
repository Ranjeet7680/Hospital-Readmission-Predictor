import os
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    # 16:9 Widescreen dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Color Palette
    BG_DARK = RGBColor(0, 47, 108)        # #002F6C Navy
    BG_LIGHT = RGBColor(248, 249, 253)    # #F8F9FD Light
    PRIMARY_BLUE = RGBColor(0, 91, 191)   # #005BBF Primary Blue
    ACCENT_CYAN = RGBColor(34, 211, 238)  # #22D3EE Cyan
    CARD_BG = RGBColor(255, 255, 255)     # #FFFFFF White
    TEXT_DARK = RGBColor(27, 27, 31)      # #1B1B1F Dark text
    TEXT_MUTED = RGBColor(91, 95, 100)    # #5B5F64 Secondary text
    ACCENT_GREEN = RGBColor(20, 108, 46)  # #146C2E Green
    ACCENT_RED = RGBColor(186, 26, 26)    # #BA1A1A Red
    ACCENT_AMBER = RGBColor(180, 100, 0)  # #B46400 Amber

    def add_header(slide, title_text, category="LUMINIX'26 HACKATHON • NEXORA TEAM"):
        # Header category banner
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.35))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category.upper()
        p_cat.font.size = Pt(11)
        p_cat.font.bold = True
        p_cat.font.color.rgb = PRIMARY_BLUE

        # Main Slide Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.75), Inches(11.7), Inches(0.7))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.size = Pt(24)
        p_title.font.bold = True
        p_title.font.color.rgb = BG_DARK

        # Top decorative line
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.5), Inches(11.7), Inches(0.04))
        line.fill.solid()
        line.fill.fore_color.rgb = PRIMARY_BLUE
        line.line.color.rgb = PRIMARY_BLUE

    def create_card(slide, left, top, width, height, bg_color=CARD_BG, border_color=RGBColor(218, 224, 233)):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color
        card.line.color.rgb = border_color
        card.line.width = Pt(1.5)
        return card

    # =========================================================================
    # SLIDE 1: TITLE SLIDE (Dark Hero Background)
    # =========================================================================
    slide1 = prs.slides.add_slide(blank_layout)
    bg1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = BG_DARK
    bg1.line.fill.background()

    # Hackathon Badge
    badge = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(1.0), Inches(4.5), Inches(0.5))
    badge.fill.solid()
    badge.fill.fore_color.rgb = RGBColor(0, 74, 158)
    badge.line.color.rgb = ACCENT_CYAN
    tf_b = badge.text_frame
    p_b = tf_b.paragraphs[0]
    p_b.text = "★ LUMINIX'26 HACKATHON PROJECT"
    p_b.font.size = Pt(11)
    p_b.font.bold = True
    p_b.font.color.rgb = ACCENT_CYAN
    p_b.alignment = PP_ALIGN.CENTER

    # Title & Subtitle
    t_box1 = slide1.shapes.add_textbox(Inches(0.9), Inches(1.7), Inches(11.5), Inches(2.2))
    tf1 = t_box1.text_frame
    tf1.word_wrap = True
    p1 = tf1.paragraphs[0]
    p1.text = "Hospital Readmission Predictor"
    p1.font.size = Pt(38)
    p1.font.bold = True
    p1.font.color.rgb = RGBColor(255, 255, 255)

    p2 = tf1.add_paragraph()
    p2.text = "AI-Powered Clinical Intelligence, Explainable XAI & Closed-Loop Care Platform"
    p2.font.size = Pt(20)
    p2.font.bold = True
    p2.font.color.rgb = ACCENT_CYAN

    p3 = tf1.add_paragraph()
    p3.text = "Predict risk. Explain insights. Connect care."
    p3.font.size = Pt(15)
    p3.font.italic = True
    p3.font.color.rgb = RGBColor(216, 226, 255)

    # Team Nexora & Info Card
    card_info = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(4.4), Inches(11.5), Inches(2.2))
    card_info.fill.solid()
    card_info.fill.fore_color.rgb = RGBColor(11, 35, 71)
    card_info.line.color.rgb = RGBColor(0, 91, 191)
    tf_info = card_info.text_frame
    tf_info.word_wrap = True

    p_team = tf_info.paragraphs[0]
    p_team.text = "TEAM: NEXORA TEAM"
    p_team.font.size = Pt(18)
    p_team.font.bold = True
    p_team.font.color.rgb = RGBColor(255, 255, 255)

    p_lead = tf_info.add_paragraph()
    p_lead.text = "Team Leader: Ranjeet Kumar   |   Email: rajranjeet7680@gmail.com"
    p_lead.font.size = Pt(14)
    p_lead.font.bold = True
    p_lead.font.color.rgb = ACCENT_CYAN

    p_core = tf_info.add_paragraph()
    p_core.text = "Core Tech: XGBoost (0.9794 AUC) • PyTorch Deep Learning • TreeSHAP XAI • PPO Reinforcement Learning • CareAI Telemedicine • Digital QR Pass"
    p_core.font.size = Pt(12)
    p_core.font.color.rgb = RGBColor(200, 215, 240)

    # =========================================================================
    # SLIDE 2: PROBLEM STATEMENT
    # =========================================================================
    slide2 = prs.slides.add_slide(blank_layout)
    add_header(slide2, "1. Problem Statement: The $26B+ Healthcare Readmission Crisis")

    cards_data2 = [
        ("The Financial & Clinical Toll", "$26+ Billion / Year", "30-day unplanned readmissions burden healthcare networks globally. In diabetic & cardiac cohorts, readmissions exceed 20%, resulting in heavy penalties under HRRP.", ACCENT_RED),
        ("Post-Discharge Blind Spots", "Fragmented Follow-up", "Clinical care teams lack automated risk triage during the critical 72-hour discharge transition window, missing acute physiological deterioration signs.", ACCENT_AMBER),
        ("Black-Box AI Distrust", "Opaque Algorithms", "Traditional AI systems output raw scores without interpretable factor attribution, preventing clinicians from understanding why a patient was flagged.", PRIMARY_BLUE)
    ]
    for i, (title, highlight, desc, col) in enumerate(cards_data2):
        x = Inches(0.8 + i * 3.95)
        create_card(slide2, x, Inches(1.8), Inches(3.8), Inches(5.0))
        tb = slide2.shapes.add_textbox(x + Inches(0.2), Inches(2.0), Inches(3.4), Inches(4.5))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = col

        p_h = tf.add_paragraph()
        p_h.text = highlight
        p_h.font.size = Pt(20)
        p_h.font.bold = True
        p_h.font.color.rgb = BG_DARK

        p_d = tf.add_paragraph()
        p_d.text = "\n" + desc
        p_d.font.size = Pt(13)
        p_d.font.color.rgb = TEXT_MUTED

    # =========================================================================
    # SLIDE 3: OUR SOLUTION
    # =========================================================================
    slide3 = prs.slides.add_slide(blank_layout)
    add_header(slide3, "2. Solution Overview: HRP Clinical Intelligence Platform")

    sol_cards = [
        ("PREDICT", "ML & Deep Learning", "Calibrated 30-day readmission risk probability using XGBoost (0.9794 AUC) & PyTorch Tabular Transformers.", PRIMARY_BLUE),
        ("UNDERSTAND", "Explainable AI (TreeSHAP)", "Patient-specific feature waterfalls explaining exact risk drivers (prior admissions, creatinine, polypharmacy).", RGBColor(0, 130, 150)),
        ("OPTIMIZE", "Reinforcement Learning", "PPO Care Pathway Digital Twin simulating personalized follow-up actions under strict clinical safety rules.", ACCENT_GREEN),
        ("CONNECT", "CareAI & Telemedicine", "Doctor video consultations with bilingual live subtitles (English ↔ हिन्दी) and QR health passes.", RGBColor(120, 50, 150))
    ]
    for i, (tag, title, desc, col) in enumerate(sol_cards):
        x = Inches(0.8 + i * 2.95)
        create_card(slide3, x, Inches(1.8), Inches(2.8), Inches(5.0))
        tb = slide3.shapes.add_textbox(x + Inches(0.15), Inches(2.0), Inches(2.5), Inches(4.5))
        tf = tb.text_frame
        tf.word_wrap = True

        p_tag = tf.paragraphs[0]
        p_tag.text = tag
        p_tag.font.size = Pt(12)
        p_tag.font.bold = True
        p_tag.font.color.rgb = col

        p_t = tf.add_paragraph()
        p_t.text = title
        p_t.font.size = Pt(16)
        p_t.font.bold = True
        p_t.font.color.rgb = BG_DARK

        p_d = tf.add_paragraph()
        p_d.text = "\n" + desc
        p_d.font.size = Pt(12.5)
        p_d.font.color.rgb = TEXT_MUTED

    # =========================================================================
    # SLIDE 4: DATASET & COHORT FOUNDATION
    # =========================================================================
    slide4 = prs.slides.add_slide(blank_layout)
    add_header(slide4, "3. Dataset & Cohort: Diabetes 130-US Hospitals (1999–2008)")

    # Left Stats (Span 4)
    create_card(slide4, Inches(0.8), Inches(1.8), Inches(5.6), Inches(5.0))
    tb_ds_l = slide4.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.2), Inches(4.5))
    tf_ds_l = tb_ds_l.text_frame
    tf_ds_l.word_wrap = True
    
    p = tf_ds_l.paragraphs[0]
    p.text = "Benchmark Clinical Scale"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = BG_DARK

    stats = [
        ("101,766 Encounters", "Diabetic inpatient hospitalizations analyzed"),
        ("130 Hospitals", "Geographically diverse US healthcare facilities"),
        ("10-Year Study Window", "Longitudinal dataset covering 1999 to 2008"),
        ("50 Clinical Features", "Demographics, admission source, labs, 23 medications")
    ]
    for num, lbl in stats:
        p_num = tf_ds_l.add_paragraph()
        p_num.text = "• " + num + ": " + lbl
        p_num.font.size = Pt(13)
        p_num.font.color.rgb = TEXT_DARK
        p_num.font.bold = True

    # Right Preprocessing (Span 4)
    create_card(slide4, Inches(6.8), Inches(1.8), Inches(5.7), Inches(5.0))
    tb_ds_r = slide4.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.3), Inches(4.5))
    tf_ds_r = tb_ds_r.text_frame
    tf_ds_r.word_wrap = True

    p_r = tf_ds_r.paragraphs[0]
    p_r.text = "10-Stage Clinical Preprocessing Pipeline"
    p_r.font.size = Pt(18)
    p_r.font.bold = True
    p_r.font.color.rgb = PRIMARY_BLUE

    steps = [
        "1. Missingness Treatment: Removed high-missing 'weight' & 'payer_code'",
        "2. Clinical Target Encoding: Binary 30-day readmission (<30 days vs NO/>30)",
        "3. Comorbidity Grouping: ICD-9 mapping (Circulatory, Endocrine, Respiratory)",
        "4. Feature Engineering: Prior utilization index, polypharmacy score, A1c anomaly",
        "5. Stratified Holdout Split: 80% Train / 20% Unseen Clinical Test Set"
    ]
    for s in steps:
        p_s = tf_ds_r.add_paragraph()
        p_s.text = s
        p_s.font.size = Pt(12)
        p_s.font.color.rgb = TEXT_MUTED

    # =========================================================================
    # SLIDE 5: ML & DEEP LEARNING ARCHITECTURE
    # =========================================================================
    slide5 = prs.slides.add_slide(blank_layout)
    add_header(slide5, "4. ML & Deep Learning: Multi-Model Evaluation")

    # Table of models
    create_card(slide5, Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
    tb_ml = slide5.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.3), Inches(4.5))
    tf_ml = tb_ml.text_frame
    tf_ml.word_wrap = True

    p = tf_ml.paragraphs[0]
    p.text = "Model Benchmark Leaderboard on Holdout Test Set"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = BG_DARK

    models_info = [
        ("XGBoost v2.4.1 (Champion)", "0.9794", "93.7%", "90.2%", "92.4%", "Optimized GBDT with TreeSHAP support"),
        ("LightGBM Classifier", "0.9712", "92.4%", "88.6%", "90.8%", "Fast histogram gradient boosting"),
        ("Random Forest (200 Trees)", "0.9645", "91.8%", "87.1%", "89.5%", "Bagged decision tree ensemble"),
        ("PyTorch Tabular Transformer", "0.9580", "90.9%", "86.4%", "88.2%", "Self-attention on tabular embeddings"),
        ("Multi-Layer Perceptron (ANN)", "0.9420", "89.5%", "84.2%", "86.8%", "Deep dense neural network with Dropout"),
        ("Logistic Regression Baseline", "0.8840", "82.1%", "76.5%", "78.9%", "L2 Regularized linear baseline")
    ]
    for name, auc, acc, sens, f1, notes in models_info:
        p_m = tf_ml.add_paragraph()
        p_m.text = f"• {name}  |  ROC-AUC: {auc}  |  Accuracy: {acc}  |  Sensitivity: {sens}  |  F1: {f1}"
        p_m.font.size = Pt(13)
        if "Champion" in name:
            p_m.font.bold = True
            p_m.font.color.rgb = PRIMARY_BLUE
        else:
            p_m.font.color.rgb = TEXT_DARK

    # =========================================================================
    # SLIDE 6: EXPLAINABLE AI (XAI) & SHAP
    # =========================================================================
    slide6 = prs.slides.add_slide(blank_layout)
    add_header(slide6, "5. Explainable AI: TreeSHAP Attribution & Transparency")

    create_card(slide6, Inches(0.8), Inches(1.8), Inches(5.6), Inches(5.0))
    tb_xai_l = slide6.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.2), Inches(4.5))
    tf_xai_l = tb_xai_l.text_frame
    tf_xai_l.word_wrap = True

    p = tf_xai_l.paragraphs[0]
    p.text = "Local Patient Factor Attribution (SHAP Waterfall)"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = BG_DARK

    factors = [
        ("Prior Inpatient Admissions (2x)", "+24.0% Risk", "History of multiple acute admissions in 90d"),
        ("Elevated Serum Creatinine (1.60)", "+16.0% Risk", "Renal impairment accelerating CHF readmission"),
        ("Polypharmacy Count (8 Meds)", "+10.2% Risk", "Complex medication regimen & drug interaction risk"),
        ("Long Length of Stay (9 Days)", "+8.5% Risk", "Severe acute inpatient hospitalization duration")
    ]
    for feat, shift, exp in factors:
        p_f = tf_xai_l.add_paragraph()
        p_f.text = f"▲ {feat} ({shift}): {exp}"
        p_f.font.size = Pt(12)
        p_f.font.color.rgb = ACCENT_RED

    create_card(slide6, Inches(6.8), Inches(1.8), Inches(5.7), Inches(5.0))
    tb_xai_r = slide6.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.3), Inches(4.5))
    tf_xai_r = tb_xai_r.text_frame
    tf_xai_r.word_wrap = True

    p_xr = tf_xai_r.paragraphs[0]
    p_xr.text = "Why Clinical Explainability Matters"
    p_xr.font.size = Pt(16)
    p_xr.font.bold = True
    p_xr.font.color.rgb = PRIMARY_BLUE

    xai_points = [
        "• Instant Physician Trust: Doctors see exact physiological drivers behind every score.",
        "• Targeted Interventions: Rather than generic care, clinicians fix specific drivers (e.g. nephrology consult for creatinine).",
        "• Counterfactual Insights: Shows clinicians what changes would reduce patient risk below 30%.",
        "• Regulatory Compliance: Satisfies FDA, EU AI Act, and HIPAA algorithmic explainability rules."
    ]
    for pt in xai_points:
        p_pt = tf_xai_r.add_paragraph()
        p_pt.text = pt
        p_pt.font.size = Pt(12.5)
        p_pt.font.color.rgb = TEXT_MUTED

    # =========================================================================
    # SLIDE 7: REINFORCEMENT LEARNING (RL)
    # =========================================================================
    slide7 = prs.slides.add_slide(blank_layout)
    add_header(slide7, "6. Reinforcement Learning: Care Pathway Optimization")

    rl_items = [
        ("6-Stage Care MDP", "States: t0 Inpatient, t1 Discharge, t2 72h Follow-up, t3 Day-7, t4 Day-14, t5 Day-30 Outcome."),
        ("PPO Agent Policy", "Trained with Proximal Policy Optimization to maximize reward (Readmission Averted: +100)."),
        ("Deterministic Safety Guardrails", "Hard safety constraints prevent unsafe suggestions. Clinician review is strictly required."),
        ("Digital Twin Simulator", "Runs 'what-if' simulations comparing Standard Care (68% risk) vs RL Pathway (26% risk).")
    ]
    for i, (title, desc) in enumerate(rl_items):
        y = Inches(1.8 + i * 1.25)
        create_card(slide7, Inches(0.8), y, Inches(11.7), Inches(1.1))
        tb = slide7.shapes.add_textbox(Inches(1.0), y + Inches(0.1), Inches(11.3), Inches(0.9))
        tf = tb.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = f"{i+1}. {title}"
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_BLUE

        p_d = tf.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(12.5)
        p_d.font.color.rgb = TEXT_DARK

    # =========================================================================
    # SLIDE 8: CAREAI & TELEMEDICINE
    # =========================================================================
    slide8 = prs.slides.add_slide(blank_layout)
    add_header(slide8, "7. CareAI Copilot & Bilingual Telemedicine")

    create_card(slide8, Inches(0.8), Inches(1.8), Inches(5.6), Inches(5.0))
    tb_c_l = slide8.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.2), Inches(4.5))
    tf_c_l = tb_c_l.text_frame
    tf_c_l.word_wrap = True

    p = tf_c_l.paragraphs[0]
    p.text = "Doctor Video Consultation & WebRTC"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = BG_DARK

    c_points = [
        "• Real-Time Video & Audio: Encrypted clinician-to-patient telemedicine.",
        "• Dual Live Subtitles: Line-by-line synchronized English & हिन्दी captioning.",
        "• Web Audio API Synthesis: Native ringtones, acoustic upload chords, click feedback.",
        "• Automated Progress Notes: Drafts clinical SOAP notes during consultation."
    ]
    for cp in c_points:
        p_cp = tf_c_l.add_paragraph()
        p_cp.text = cp
        p_cp.font.size = Pt(12.5)
        p_cp.font.color.rgb = TEXT_MUTED

    create_card(slide8, Inches(6.8), Inches(1.8), Inches(5.7), Inches(5.0))
    tb_c_r = slide8.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.3), Inches(4.5))
    tf_c_r = tb_c_r.text_frame
    tf_c_r.word_wrap = True

    p_r = tf_c_r.paragraphs[0]
    p_r.text = "CareAI Conversational Healthcare Copilot"
    p_r.font.size = Pt(16)
    p_r.font.bold = True
    p_r.font.color.rgb = PRIMARY_BLUE

    ai_caps = [
        "• Risk Explanation in Hindi: Translates complex SHAP telemetry into conversational हिन्दी.",
        "• Medical Document Q&A: Ingests lab reports and cites exact page/biomarker sources.",
        "• Discharge Instructions: Prepares patient-friendly medication checklists.",
        "• Human Oversight Disclaimer: Always marks outputs as assistive suggestions."
    ]
    for ac in ai_caps:
        p_ac = tf_c_r.add_paragraph()
        p_ac.text = ac
        p_ac.font.size = Pt(12.5)
        p_ac.font.color.rgb = TEXT_MUTED

    # =========================================================================
    # SLIDE 9: MEDICAL REPORT OCR & CERTIFICATES
    # =========================================================================
    slide9 = prs.slides.add_slide(blank_layout)
    add_header(slide9, "8. Medical Documents: OCR & Digital Certificates")

    doc_steps = [
        ("1. PDF Ingestion", "Upload laboratory panels, CBCs, and hospital discharge summaries."),
        ("2. Structured OCR Extraction", "Extracts Serum Creatinine, BUN, Hemoglobin, HbA1c, and reference ranges."),
        ("3. AI Biomarker Anomaly Detection", "Flags high/low values with visual severity tags (e.g. Creatinine 1.60 mg/dL HIGH)."),
        ("4. Medical Leave Certificates", "Generates verifiable convalescence certificates with doctor digital signatures.")
    ]
    for i, (title, desc) in enumerate(doc_steps):
        x = Inches(0.8 + i * 2.95)
        create_card(slide9, x, Inches(1.8), Inches(2.8), Inches(5.0))
        tb = slide9.shapes.add_textbox(x + Inches(0.15), Inches(2.0), Inches(2.5), Inches(4.5))
        tf = tb.text_frame
        tf.word_wrap = True

        p_t = tf.paragraphs[0]
        p_t.text = title
        p_t.font.size = Pt(15)
        p_t.font.bold = True
        p_t.font.color.rgb = PRIMARY_BLUE

        p_d = tf.add_paragraph()
        p_d.text = "\n" + desc
        p_d.font.size = Pt(12.5)
        p_d.font.color.rgb = TEXT_MUTED

    # =========================================================================
    # SLIDE 10: DIGITAL HEALTH ID & QR SYSTEM
    # =========================================================================
    slide10 = prs.slides.add_slide(blank_layout)
    add_header(slide10, "9. Digital Health ID & Working QR Engine")

    create_card(slide10, Inches(0.8), Inches(1.8), Inches(5.6), Inches(5.0))
    tb_id_l = slide10.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.2), Inches(4.5))
    tf_id_l = tb_id_l.text_frame
    tf_id_l.word_wrap = True

    p = tf_id_l.paragraphs[0]
    p.text = "Interactive 3D Digital Health ID Card"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = BG_DARK

    id_feats = [
        "• 3D Flip Perspective: Front identity face + back verification face.",
        "• Verified Credential: Eleanor Vance (#HRP-2026-0001042), Level 3 verified.",
        "• Pure SVG QR Generator: Scalable vector QR without external internet dependencies.",
        "• Instant Lost ID Reporting: One-click token invalidation & key regeneration."
    ]
    for f in id_feats:
        p_f = tf_id_l.add_paragraph()
        p_f.text = f
        p_f.font.size = Pt(12.5)
        p_f.font.color.rgb = TEXT_MUTED

    create_card(slide10, Inches(6.8), Inches(1.8), Inches(5.7), Inches(5.0))
    tb_id_r = slide10.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.3), Inches(4.5))
    tf_id_r = tb_id_r.text_frame
    tf_id_r.word_wrap = True

    p_r = tf_id_r.paragraphs[0]
    p_r.text = "Multi-Purpose Passes & Scanner"
    p_r.font.size = Pt(16)
    p_r.font.bold = True
    p_r.font.color.rgb = PRIMARY_BLUE

    qr_types = [
        "• In-Browser Camera Scanner: Animated laser scanner with instant verification.",
        "• Doctor Profile QR: Public check endpoint for clinician board certification.",
        "• Appointment Pass: Fast-track clinic check-in terminal pass.",
        "• Temporary Document Sharing: Auto-expiring access links (1h, 24h, 7d) with instant revocation."
    ]
    for qt in qr_types:
        p_qt = tf_id_r.add_paragraph()
        p_qt.text = qt
        p_qt.font.size = Pt(12.5)
        p_qt.font.color.rgb = TEXT_MUTED

    # =========================================================================
    # SLIDE 11: SECURITY, PRIVACY & RBAC
    # =========================================================================
    slide11 = prs.slides.add_slide(blank_layout)
    add_header(slide11, "10. Security & HIPAA-Aligned Governance")

    sec_cards = [
        ("4-Tier RBAC", "Strict isolation between Patient, Doctor, Care Coordinator, and Administrator."),
        ("Multi-Factor Authentication", "6-Digit Time-Based OTP (TOTP) and WebAuthn / FIDO2 Passkeys."),
        ("Break-Glass Protocol", "Emergency access overrides with real-time audit alerts & chief medical officer logging."),
        ("HIPAA Data Portability", "One-click 'Download My Data' personal health JSON archive export.")
    ]
    for i, (title, desc) in enumerate(sec_cards):
        y = Inches(1.8 + i * 1.25)
        create_card(slide11, Inches(0.8), y, Inches(11.7), Inches(1.1))
        tb = slide11.shapes.add_textbox(Inches(1.0), y + Inches(0.1), Inches(11.3), Inches(0.9))
        tf = tb.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = f"🔒 {title}"
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_BLUE

        p_d = tf.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(12.5)
        p_d.font.color.rgb = TEXT_DARK

    # =========================================================================
    # SLIDE 12: SYSTEM ARCHITECTURE & TECH STACK
    # =========================================================================
    slide12 = prs.slides.add_slide(blank_layout)
    add_header(slide12, "11. Full-Stack System Architecture")

    arch_layers = [
        ("Frontend Presentation", "Google Material 3, Tailwind CSS, Responsive Breakpoints, Web Audio API, i18n English/Hindi engine."),
        ("Application Backend", "FastAPI (Async Python 3.11), Jinja2 Templates, Secure Session Store, Token Verification Engine."),
        ("AI & Decision Core", "PyTorch 2.4, Scikit-Learn, XGBoost v2.4.1, LightGBM, TreeSHAP, Stable-Baselines3 (PPO RL)."),
        ("Security & Documents", "WebAuthn / Passkeys, Pure SVG QR Generator, OCR Engine, PDF Report Exporter.")
    ]
    for i, (layer, tech) in enumerate(arch_layers):
        y = Inches(1.8 + i * 1.25)
        create_card(slide12, Inches(0.8), y, Inches(11.7), Inches(1.1))
        tb = slide12.shapes.add_textbox(Inches(1.0), y + Inches(0.1), Inches(11.3), Inches(0.9))
        tf = tb.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = f"⚙️ {layer}"
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = BG_DARK

        p_d = tf.add_paragraph()
        p_d.text = tech
        p_d.font.size = Pt(12.5)
        p_d.font.color.rgb = PRIMARY_BLUE

    # =========================================================================
    # SLIDE 13: TEST VERIFICATION & VALIDATION
    # =========================================================================
    slide13 = prs.slides.add_slide(blank_layout)
    add_header(slide13, "12. Live Validation & Automated Test Results")

    create_card(slide13, Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
    tb_val = slide13.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.3), Inches(4.5))
    tf_val = tb_val.text_frame
    tf_val.word_wrap = True

    p = tf_val.paragraphs[0]
    p.text = "Automated Test Suite Status: 18 / 18 Tests Passing (100% Success)"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GREEN

    tests_list = [
        ("✓ Model Inference Accuracy", "Evaluates high-risk vs low-risk patient payloads with calibrated scores."),
        ("✓ Authentication & MFA Security", "Tests password hashing, 6-digit OTP verification, and RBAC authorization."),
        ("✓ Break-Glass Protocol", "Simulates emergency access override with audit logging."),
        ("✓ Medical Document & Lab OCR", "Validates structured biomarker extraction and anomaly matching."),
        ("✓ Reinforcement Learning Policy", "Verifies PPO care journey transitions and hard safety constraint enforcement."),
        ("✓ QR Token Lifecycle & Rotation", "Validates generation, public check, temporary sharing, and lost ID invalidation."),
        ("✓ Full Web Route Smoke Suite", "Validates HTTP 200 across 55+ clinical, AI, and portal endpoints.")
    ]
    for t_name, t_desc in tests_list:
        p_t = tf_val.add_paragraph()
        p_t.text = f"{t_name}: {t_desc}"
        p_t.font.size = Pt(12.5)
        p_t.font.color.rgb = TEXT_DARK

    # =========================================================================
    # SLIDE 14: BUSINESS IMPACT & ROI
    # =========================================================================
    slide14 = prs.slides.add_slide(blank_layout)
    add_header(slide14, "13. Clinical Impact & Healthcare ROI")

    roi_cards = [
        ("18.6% Reduction", "In preventable 30-day readmissions through proactive 72h follow-ups.", ACCENT_GREEN),
        ("HRRP Penalty Avoidance", "Protects hospital Medicare reimbursement margins by keeping readmission rates low.", PRIMARY_BLUE),
        ("4.5h Faster Triage", "Accelerates clinical review with automated TreeSHAP driver attribution.", ACCENT_AMBER),
        ("Higher Patient Trust", "Bilingual English/Hindi patient engagement and digital health passes.", RGBColor(120, 50, 150))
    ]
    for i, (highlight, desc, col) in enumerate(roi_cards):
        x = Inches(0.8 + i * 2.95)
        create_card(slide14, x, Inches(1.8), Inches(2.8), Inches(5.0))
        tb = slide14.shapes.add_textbox(x + Inches(0.15), Inches(2.0), Inches(2.5), Inches(4.5))
        tf = tb.text_frame
        tf.word_wrap = True

        p_h = tf.paragraphs[0]
        p_h.text = highlight
        p_h.font.size = Pt(18)
        p_h.font.bold = True
        p_h.font.color.rgb = col

        p_d = tf.add_paragraph()
        p_d.text = "\n" + desc
        p_d.font.size = Pt(13)
        p_d.font.color.rgb = TEXT_MUTED

    # =========================================================================
    # SLIDE 15: CONCLUSION & FUTURE ROADMAP (Dark Theme)
    # =========================================================================
    slide15 = prs.slides.add_slide(blank_layout)
    bg15 = slide15.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg15.fill.solid()
    bg15.fill.fore_color.rgb = BG_DARK
    bg15.line.fill.background()

    tb15 = slide15.shapes.add_textbox(Inches(0.9), Inches(1.0), Inches(11.5), Inches(5.5))
    tf15 = tb15.text_frame
    tf15.word_wrap = True

    p = tf15.paragraphs[0]
    p.text = "14. Future Roadmap & Hackathon Conclusion"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    roadmap = [
        ("Phase 1: EHR FHIR Integration", "Direct HL7 FHIR bidirectional pipeline with Epic & Cerner hospital databases."),
        ("Phase 2: Continuous Wearable IoT", "Smartwatch vitals & continuous glucose monitoring (CGM) telemetry stream."),
        ("Phase 3: Federated Clinical Learning", "Privacy-preserving multi-hospital collaborative model updates.")
    ]
    for title, desc in roadmap:
        p_r = tf15.add_paragraph()
        p_r.text = f"🚀 {title}: {desc}"
        p_r.font.size = Pt(14)
        p_r.font.bold = True
        p_r.font.color.rgb = ACCENT_CYAN

    p_thx = tf15.add_paragraph()
    p_thx.text = "\nThank You!  |  LUMINIX'26 Presentation"
    p_thx.font.size = Pt(24)
    p_thx.font.bold = True
    p_thx.font.color.rgb = RGBColor(255, 255, 255)

    p_cred = tf15.add_paragraph()
    p_cred.text = "Team: Nexora Team   |   Leader: Ranjeet Kumar (rajranjeet7680@gmail.com)\nGitHub: https://github.com/Ranjeet7680/Hospital-Readmission-Predictor"
    p_cred.font.size = Pt(14)
    p_cred.font.color.rgb = RGBColor(200, 220, 255)

    output_path = os.path.join(os.getcwd(), "Hospital_Readmission_Predictor_LUMINIX26.pptx")
    prs.save(output_path)
    print(f"Presentation saved successfully to {output_path}")

if __name__ == "__main__":
    create_presentation()
